#!/usr/bin/env python3
"""Approximate PNG renderer for a python-pptx deck (layout QA only).
Usage: render_deck.py deck.pptx outdir [dpi]
"""
import sys, os, math, glob
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Emu
from pptx.oxml.ns import qn
from pptx.shapes.connector import Connector
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

PPTX, OUT = sys.argv[1], sys.argv[2]
DPI = float(sys.argv[3]) if len(sys.argv) > 3 else 80.0
os.makedirs(OUT, exist_ok=True)
EMU_PER_IN = 914400.0


def px(emu):
    return emu / EMU_PER_IN * DPI


SUP = "/System/Library/Fonts/Supplemental/"
FONT_FILES = {
    ("Georgia", False, False): SUP + "Georgia.ttf", ("Georgia", True, False): SUP + "Georgia Bold.ttf",
    ("Georgia", False, True): SUP + "Georgia Italic.ttf", ("Georgia", True, True): SUP + "Georgia Bold Italic.ttf",
    ("Calibri", False, False): SUP + "Arial.ttf", ("Calibri", True, False): SUP + "Arial Bold.ttf",
    ("Calibri", False, True): SUP + "Arial Italic.ttf", ("Calibri", True, True): SUP + "Arial Bold Italic.ttf",
    ("Consolas", False, False): "/System/Library/Fonts/Menlo.ttc", ("Consolas", True, False): "/System/Library/Fonts/Menlo.ttc",
    ("Consolas", False, True): "/System/Library/Fonts/Menlo.ttc", ("Consolas", True, True): "/System/Library/Fonts/Menlo.ttc",
}
_cache = {}


def font_for(name, size_pt, bold, italic):
    key = (name or "Calibri", bool(bold), bool(italic), round(size_pt, 1))
    if key in _cache:
        return _cache[key]
    path = FONT_FILES.get((key[0], key[1], key[2])) or FONT_FILES[("Calibri", key[1], key[2])]
    pxsize = max(4, int(round(size_pt * DPI / 72.0)))
    try:
        f = ImageFont.truetype(path, pxsize, index=1 if (path.endswith(".ttc") and bold) else 0)
    except Exception:
        f = ImageFont.truetype(SUP + "Arial.ttf", pxsize)
    _cache[key] = f
    return f


def color_of(rgb, default=(245, 245, 240)):
    try:
        h = str(rgb)
        return tuple(int(h[i:i + 2], 16) for i in (0, 2, 4))
    except Exception:
        return default


def para_tokens(p):
    toks = []
    for r in p.runs:
        size = r.font.size.pt if r.font.size else 14
        f = font_for(r.font.name, size, r.font.bold, r.font.italic)
        try:
            col = color_of(r.font.color.rgb)
        except Exception:
            col = (245, 245, 240)
        parts = r.text.split(" ")
        for i, w in enumerate(parts):
            toks.append((w + (" " if i < len(parts) - 1 else ""), f, col, size))
    return toks


def layout_paragraph(p, width_px):
    toks = para_tokens(p)
    sizes = [t[3] for t in toks] or [14]
    maxsize = max(sizes)
    ls = p.line_spacing if isinstance(p.line_spacing, float) else 1.1
    line_h = maxsize * DPI / 72.0 * 1.2 * ls
    after = (p.space_after.pt if p.space_after is not None else 0) * DPI / 72.0
    pPr = p._p.pPr
    marL = int(pPr.get("marL", 0)) if pPr is not None else 0
    bullet = None
    if pPr is not None:
        bc = pPr.find(qn("a:buChar"))
        if bc is not None:
            bullet = bc.get("char")
    marL_px = px(marL)
    avail = max(10, width_px - marL_px)
    lines, cur, cur_w = [], [], 0.0
    for tok in toks:
        w = tok[1].getlength(tok[0])
        if cur and cur_w + tok[1].getlength(tok[0].rstrip()) > avail:
            lines.append(cur)
            cur, cur_w = [], 0.0
        cur.append(tok)
        cur_w += w
    if cur or not lines:
        lines.append(cur)
    return lines, line_h, after, marL_px, bullet


def draw_textframe(draw, tf, x, y, w, h, flag):
    ml, mr = px(tf.margin_left), px(tf.margin_right)
    mt, mb = px(tf.margin_top), px(tf.margin_bottom)
    ix, iw = x + ml, w - ml - mr
    blocks, total = [], 0.0
    for p in tf.paragraphs:
        lines, lh, after, marL, bullet = layout_paragraph(p, iw)
        blocks.append((p, lines, lh, after, marL, bullet))
        total += len(lines) * lh + after
    anchor = tf.vertical_anchor
    inner_h = h - mt - mb
    if anchor == MSO_ANCHOR.MIDDLE:
        ty = y + mt + max(0, (inner_h - total) / 2)
    elif anchor == MSO_ANCHOR.BOTTOM:
        ty = y + mt + max(0, inner_h - total)
    else:
        ty = y + mt
    overflow = total > inner_h + 2
    cy = ty
    for p, lines, lh, after, marL, bullet in blocks:
        for li, line in enumerate(lines):
            lw = sum(t[1].getlength(t[0]) for t in line)
            if p.alignment == PP_ALIGN.CENTER:
                lx = ix + marL + (iw - marL - lw) / 2
            elif p.alignment == PP_ALIGN.RIGHT:
                lx = ix + iw - lw
            else:
                lx = ix + marL
            if bullet and li == 0 and line and any(t[0].strip() for t in line):
                draw.text((ix + max(0, marL - px(182880)), cy), bullet, font=line[0][1], fill=(34, 211, 238))
            for tok in line:
                draw.text((lx, cy), tok[0], font=tok[1], fill=tok[2])
                lx += tok[1].getlength(tok[0])
            cy += lh
        cy += after
    if overflow:
        flag.append(f"text overflow by {(total - inner_h) / DPI:.2f}in: {tf.text[:50]!r}")
        draw.rectangle([x, y, x + w, y + h], outline=(255, 60, 60), width=2)


def shape_fill(shape):
    try:
        if shape.fill.type == 1:
            return color_of(shape.fill.fore_color.rgb)
    except Exception:
        pass
    return None


def shape_line(shape):
    try:
        if shape.line.fill.type == 1:
            return color_of(shape.line.color.rgb), max(1, int(round(px(shape.line.width) if shape.line.width else 1)))
    except Exception:
        pass
    return None, 0


def render_slide(slide, W, H):
    img = Image.new("RGB", (int(W), int(H)), (20, 23, 28))
    draw = ImageDraw.Draw(img)
    flags = []
    for shape in slide.shapes:
        x, y, w, h = px(shape.left), px(shape.top), px(shape.width), px(shape.height)
        if isinstance(shape, Connector):
            xfrm = shape._element.find(".//" + qn("a:xfrm"))
            fh = xfrm is not None and xfrm.get("flipH") == "1"
            fv = xfrm is not None and xfrm.get("flipV") == "1"
            x1, x2 = (x + w, x) if fh else (x, x + w)
            y1, y2 = (y + h, y) if fv else (y, y + h)
            col, lw = shape_line(shape)
            draw.line([x1, y1, x2, y2], fill=col or (100, 100, 100), width=max(1, lw))
            ln = shape._element.find(".//" + qn("a:ln"))
            if ln is not None and ln.find(qn("a:tailEnd")) is not None:
                ang = math.atan2(y2 - y1, x2 - x1)
                for d in (0.5, -0.5):
                    draw.line([x2, y2, x2 - 9 * math.cos(ang + d), y2 - 9 * math.sin(ang + d)], fill=col or (100, 100, 100), width=max(1, lw))
            if ln is not None and ln.find(qn("a:headEnd")) is not None:
                ang = math.atan2(y1 - y2, x1 - x2)
                for d in (0.5, -0.5):
                    draw.line([x1, y1, x1 - 9 * math.cos(ang + d), y1 - 9 * math.sin(ang + d)], fill=col or (100, 100, 100), width=max(1, lw))
            continue
        if shape.has_table:
            tbl = shape.table
            cw = [px(c.width) for c in tbl.columns]
            cy = y
            for r in tbl.rows:
                rh = px(r.height)
                need = 0
                for ci, cell in enumerate(r.cells):
                    iw = cw[ci] - px(cell.margin_left) - px(cell.margin_right)
                    tot = 0
                    for p in cell.text_frame.paragraphs:
                        lines, lh, after, marL, bullet = layout_paragraph(p, iw)
                        tot += len(lines) * lh + after
                    need = max(need, tot + px(cell.margin_top) + px(cell.margin_bottom) + 2)
                rh = max(rh, need)
                cx = x
                for ci, cell in enumerate(r.cells):
                    fill = None
                    try:
                        if cell.fill.type == 1:
                            fill = color_of(cell.fill.fore_color.rgb)
                    except Exception:
                        pass
                    draw.rectangle([cx, cy, cx + cw[ci], cy + rh], fill=fill or (35, 41, 50), outline=(58, 65, 75), width=1)
                    draw_textframe(draw, cell.text_frame, cx, cy, cw[ci], rh, [])
                    cx += cw[ci]
                cy += rh
            if cy > H - 2:
                flags.append(f"table runs past slide bottom by {(cy - H) / DPI:.2f}in")
                draw.rectangle([x, y, x + w, min(cy, H - 1)], outline=(255, 60, 60), width=3)
            elif cy > y + h + 2:
                flags.append(f"table grew {(cy - y - h) / DPI:.2f}in below its frame (bottom at {cy / DPI:.2f}in)")
            continue
        fill = shape_fill(shape)
        lcol, lw = shape_line(shape)
        try:
            ast = shape.auto_shape_type
        except Exception:
            ast = None
        if fill or lcol:
            if ast == MSO_SHAPE.OVAL:
                draw.ellipse([x, y, x + w, y + h], fill=fill, outline=lcol, width=lw)
            elif ast == MSO_SHAPE.CHEVRON:
                d = min(w, h) * 0.25
                pts = [(x, y), (x + w - d, y), (x + w, y + h / 2), (x + w - d, y + h), (x, y + h), (x + d, y + h / 2)]
                draw.polygon(pts, fill=fill, outline=lcol)
            elif ast == MSO_SHAPE.ROUNDED_RECTANGLE:
                rad = min(w, h) * (shape.adjustments[0] if len(shape.adjustments) else 0.1)
                draw.rounded_rectangle([x, y, x + w, y + h], radius=rad, fill=fill, outline=lcol, width=lw)
            else:
                draw.rectangle([x, y, x + w, y + h], fill=fill, outline=lcol, width=lw)
        if shape.has_text_frame and shape.text_frame.text.strip():
            draw_textframe(draw, shape.text_frame, x, y, w, h, flags)
    return img, flags


prs = Presentation(PPTX)
W, H = px(prs.slide_width), px(prs.slide_height)
report = []
for i, s in enumerate(prs.slides, 1):
    img, flags = render_slide(s, W, H)
    img.save(f"{OUT}/slide-{i:02d}.png")
    report += [f"slide {i}: {f}" for f in flags]
files = sorted(glob.glob(f"{OUT}/slide-*.png"))
for si in range(0, len(files), 4):
    imgs = [Image.open(f) for f in files[si:si + 4]]
    w, h = imgs[0].size
    sheet = Image.new("RGB", (w * 2 + 24, h * 2 + 24), (255, 255, 255))
    for k, im in enumerate(imgs):
        sheet.paste(im, ((k % 2) * (w + 24), (k // 2) * (h + 24)))
    sheet.save(f"{OUT}/sheet-{si // 4 + 1:02d}.png")
print("\n".join(report) if report else "no overflow flags")
print(f"rendered {len(files)} slides to {OUT}")
