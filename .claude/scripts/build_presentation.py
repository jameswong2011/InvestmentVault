#!/usr/bin/env python3
"""Build the "A Second Brain for Investing" deck (self-contained; python-pptx).
Usage: python3 .claude/scripts/build_presentation.py OUT.pptx OUT_OUTLINE.md
"""

import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.oxml.ns import qn
from lxml import etree

OUT_PPTX, OUT_MD = sys.argv[1], sys.argv[2]

W, H = 13.333, 7.5
HEAD, BODY, MONO = "Georgia", "Calibri", "Consolas"
C = {
    "bg": "14171C", "panel": "1B2027", "panel2": "232932", "border": "3A414B",
    "text": "F5F5F0", "text2": "D1D5DB", "muted": "9CA3AF", "dim": "6B7280",
    "accent": "0891B2", "accent2": "22D3EE", "accentdark": "0E4F63",
    "amber": "F59E0B", "green": "34D399", "red": "F87171", "white": "FFFFFF",
}
DECK_LABEL = "Laniakea Partners · A Second Brain for Investing · August 2026"


def rgb(key):
    return RGBColor.from_string(C.get(key, key))


prs = Presentation()
prs.slide_width, prs.slide_height = Inches(W), Inches(H)
BLANK = prs.slide_layouts[6]
SLIDES = []


def _style_run(r, size, color, font, bold, italic):
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = rgb(color)


def _bullet(p, level=0, char="•", color="accent2"):
    pPr = p._p.get_or_add_pPr()
    pPr.set("marL", str(int(Inches(0.24 + 0.26 * level))))
    pPr.set("indent", str(int(-Inches(0.2))))
    buClr = etree.SubElement(pPr, qn("a:buClr"))
    s = etree.SubElement(buClr, qn("a:srgbClr"))
    s.set("val", C[color])
    buFont = etree.SubElement(pPr, qn("a:buFont"))
    buFont.set("typeface", "Arial")
    buChar = etree.SubElement(pPr, qn("a:buChar"))
    buChar.set("char", char)


ALIGN = {"l": PP_ALIGN.LEFT, "c": PP_ALIGN.CENTER, "r": PP_ALIGN.RIGHT}
ANCHOR = {"t": MSO_ANCHOR.TOP, "m": MSO_ANCHOR.MIDDLE, "b": MSO_ANCHOR.BOTTOM}


def fill_tf(tf, content, size=14, color="text", font=BODY, bold=False, italic=False,
            align="l", spacing=1.12, after=4):
    tf.word_wrap = True
    items = content if isinstance(content, list) else [content]
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        d = item if isinstance(item, dict) else {"t": item}
        p.alignment = ALIGN[d.get("align", align)]
        p.line_spacing = d.get("spacing", spacing)
        p.space_after = Pt(d.get("after", after))
        if "before" in d:
            p.space_before = Pt(d["before"])
        if d.get("bullet"):
            _bullet(p, d.get("level", 0), d.get("char", "•"), d.get("bcolor", "accent2"))
        runs = d.get("t", "")
        if isinstance(runs, str):
            runs = [(runs, {})]
        for text, st in runs:
            r = p.add_run()
            r.text = text
            _style_run(r, st.get("size", d.get("size", size)), st.get("color", d.get("color", color)),
                       st.get("font", d.get("font", font)), st.get("bold", d.get("bold", bold)),
                       st.get("italic", d.get("italic", italic)))
    return tf


def add_text(slide, x, y, w, h, content, size=14, color="text", font=BODY, bold=False, italic=False,
             align="l", anchor="t", spacing=1.12, after=4, margin=0.05, name=None):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    if name:
        tb.name = name
    tf = tb.text_frame
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = tf.margin_right = Inches(margin)
    tf.margin_top = tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = ANCHOR[anchor]
    fill_tf(tf, content, size, color, font, bold, italic, align, spacing, after)
    return tb


def add_box(slide, x, y, w, h, fill="panel", line="border", text=None, size=13, color="text", font=BODY,
            bold=False, align="c", anchor="m", shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.12,
            line_w=1.0, margin=0.12, spacing=1.1, after=3, italic=False):
    shp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        shp.adjustments[0] = radius
    if fill:
        shp.fill.solid()
        shp.fill.fore_color.rgb = rgb(fill)
    else:
        shp.fill.background()
    if line:
        shp.line.color.rgb = rgb(line)
        shp.line.width = Pt(line_w)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = tf.margin_right = Inches(margin)
    tf.margin_top = tf.margin_bottom = Inches(margin * 0.7)
    tf.vertical_anchor = ANCHOR[anchor]
    if text is not None:
        fill_tf(tf, text, size, color, font, bold, italic, align, spacing, after)
    return shp


def add_line(slide, x1, y1, x2, y2, color="border", width=1.25, head=False, tail=False, dash=False):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = rgb(color)
    conn.line.width = Pt(width)
    ln = conn.line._get_or_add_ln()
    if dash:
        pd = etree.SubElement(ln, qn("a:prstDash"))
        pd.set("val", "dash")
    if tail:
        he = etree.SubElement(ln, qn("a:headEnd"))
        he.set("type", "triangle"); he.set("w", "med"); he.set("len", "med")
    if head:
        te = etree.SubElement(ln, qn("a:tailEnd"))
        te.set("type", "triangle"); te.set("w", "med"); te.set("len", "med")
    return conn


def _cell_border(cell, color="border", w_pt=0.75):
    tcPr = cell._tc.get_or_add_tcPr()
    for i, tag in enumerate(["a:lnL", "a:lnR", "a:lnT", "a:lnB"]):
        ln = etree.Element(qn(tag))
        ln.set("w", str(int(Pt(w_pt)))); ln.set("cap", "flat"); ln.set("cmpd", "sng"); ln.set("algn", "ctr")
        sf = etree.SubElement(ln, qn("a:solidFill"))
        c = etree.SubElement(sf, qn("a:srgbClr")); c.set("val", C[color])
        pd = etree.SubElement(ln, qn("a:prstDash")); pd.set("val", "solid")
        tcPr.insert(i, ln)


def add_table(slide, x, y, w, h, data, col_widths=None, size=11.5, header=True, first_col_bold=True,
              header_size=None, anchor="t"):
    rows, cols = len(data), len(data[0])
    gs = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(h))
    tbl = gs.table
    tbl.first_row = header
    tbl.horz_banding = False
    tbl.vert_banding = False
    if col_widths:
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = Inches(cw)
    for r in range(rows):
        tbl.rows[r].height = Inches(h / rows)
        for c in range(cols):
            cell = tbl.cell(r, c)
            cell.margin_left = cell.margin_right = Inches(0.09)
            cell.margin_top = cell.margin_bottom = Inches(0.05)
            cell.vertical_anchor = ANCHOR[anchor]
            is_head = header and r == 0
            cell.fill.solid()
            cell.fill.fore_color.rgb = rgb("accentdark" if is_head else ("panel" if r % 2 else "panel2"))
            _cell_border(cell)
            val = data[r][c]
            tf = cell.text_frame
            fill_tf(tf, val, size=(header_size or size) if is_head else size,
                    color="white" if is_head else ("text" if (c == 0 and first_col_bold) else "text2"),
                    bold=is_head or (c == 0 and first_col_bold), after=2, spacing=1.08)
    return gs


def new_slide(section):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = rgb("bg")
    SLIDES.append((s, section))
    return s


def add_title(slide, title, kicker=None, title_size=28):
    if kicker:
        add_text(slide, 0.6, 0.32, 12.0, 0.35, kicker.upper(), size=11, color="accent2", font=BODY, bold=True)
    add_text(slide, 0.6, 0.62, 12.1, 0.95, title, size=title_size, color="text", font=HEAD, anchor="t", spacing=1.0)
    add_line(slide, 0.65, 1.52, 1.65, 1.52, color="accent", width=2.5)


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def divider(section, number, title, bullets, note):
    s = new_slide(section)
    add_text(s, 0.9, 1.6, 3.0, 1.4, f"0{number}", size=64, color="accent", font=HEAD)
    add_text(s, 0.9, 2.95, 11.5, 1.25, title, size=32, color="text", font=HEAD, spacing=1.0)
    add_line(s, 0.95, 4.25, 2.4, 4.25, color="accent", width=2.5)
    add_text(s, 0.9, 4.45, 11.0, 2.0, [{"t": b, "bullet": True, "size": 16, "color": "text2", "after": 6} for b in bullets])
    notes(s, note)
    return s


def stat_tile(slide, x, y, w, h, number, caption, source=None, number_color="accent2"):
    add_box(slide, x, y, w, h, fill="panel", line="border")
    add_text(slide, x + 0.15, y + 0.12, w - 0.3, 0.9, number, size=34, color=number_color, font=HEAD)
    add_text(slide, x + 0.15, y + 1.0, w - 0.3, h - 1.15, [{"t": caption, "size": 12.5, "color": "text2", "after": 4}] +
             ([{"t": source, "size": 9.5, "color": "dim"}] if source else []), spacing=1.1)


def two_panels(slide, left_title, left_items, right_title, right_items, y=1.8, h=4.9, size=13, gap=0.3,
               left_color="green", right_color="red"):
    w = (12.1 - gap) / 2
    for i, (t, items, col) in enumerate([(left_title, left_items, left_color), (right_title, right_items, right_color)]):
        x = 0.6 + i * (w + gap)
        add_box(slide, x, y, w, h, fill="panel", line="border")
        add_text(slide, x + 0.2, y + 0.15, w - 0.4, 0.45, t, size=15, color=col, font=BODY, bold=True)
        add_text(slide, x + 0.2, y + 0.65, w - 0.4, h - 0.8,
                 [{"t": it, "bullet": True, "size": size, "color": "text2", "after": 5, "bcolor": col} for it in items])





def sat(slide, x, y, w, h, title, body, fill="panel", line="border"):
    return add_box(slide, x, y, w, h, fill=fill, line=line, anchor="m", align="c", text=[
        {"t": title, "size": 13, "color": "text", "bold": True, "after": 2},
        {"t": body, "size": 10.5, "color": "muted"},
    ])



# ================================================================= SLIDES")[0]
EXPORT = SRC.split("# ---------------------------------------------------------------- footers")[1]
exec(HELPERS)


def sat(slide, x, y, w, h, title, body, fill="panel", line="border"):
    return add_box(slide, x, y, w, h, fill=fill, line=line, anchor="m", align="c", text=[
        {"t": title, "size": 13, "color": "text", "bold": True, "after": 2},
        {"t": body, "size": 10.5, "color": "muted"},
    ])



# ---------------------------------------------------------------- footers
total = len(SLIDES)
for i, (sl, section) in enumerate(SLIDES, 1):
    if i in (1, total):
        continue
    add_text(sl, 0.6, 7.05, 8.5, 0.3, DECK_LABEL + (f"  ·  {section}" if section else ""), size=9, color="dim", name="footer")
    add_text(sl, 11.3, 7.05, 1.45, 0.3, f"{i} / {total}", size=9, color="dim", align="r", name="footer")

prs.save(OUT_PPTX)


def shape_text(sh):
    out = []
    if sh.has_text_frame:
        t = "\n".join(p.text for p in sh.text_frame.paragraphs if p.text.strip())
        if t.strip():
            out.append(t)
    if sh.has_table:
        tbl = sh.table
        rows = [[c.text.replace("\n", " ").strip() for c in r.cells] for r in tbl.rows]
        out.append("| " + " | ".join(rows[0]) + " |")
        out.append("|" + "---|" * len(rows[0]))
        for r in rows[1:]:
            out.append("| " + " | ".join(r) + " |")
    return out


lines = ["---", "publish: false", "date: 2026-08-25", "tags: [meta, presentation, second-brain]",
         "status: active", "source: generated from Build documents + Website essays + vault state",
         "---", "", "# A Second Brain for Investing — presentation outline and speaker notes", "",
         f"Companion to `Build documents/Second Brain for Investing - Presentation (2026-08-25).pptx` ({total} slides, 16:9). "
         "Each section below is one slide: on-slide text first, speaker notes second. Usable as an editing outline, or as the "
         "prompt pack for regenerating the deck in another tool.", ""]
for i, (sl, section) in enumerate(SLIDES, 1):
    title = None
    body = []
    for sh in sl.shapes:
        if sh.name == "footer":
            continue
        for t in shape_text(sh):
            if title is None and sh.has_text_frame and sh.text_frame.paragraphs[0].runs and \
               sh.text_frame.paragraphs[0].runs[0].font.name == HEAD and sh.text_frame.paragraphs[0].runs[0].font.size >= Pt(26) \
               and not t.strip().isdigit():
                title = t.split("\n")[0]
            else:
                body.append(t)
    lines.append(f"## Slide {i} — {title or '(untitled)'}" + (f"  `{section}`" if section else ""))
    lines.append("")
    for b in body:
        for ln_ in b.split("\n"):
            lines.append(ln_ if ln_.startswith("|") else f"- {ln_}")
    lines.append("")
    nt = sl.notes_slide.notes_text_frame.text.strip() if sl.has_notes_slide else ""
    if nt:
        lines.append(f"> **Speaker notes:** {nt}")
        lines.append("")
with open(OUT_MD, "w", encoding="utf-8") as f:
    f.write("\n".join(lines))
print(f"wrote {OUT_PPTX} ({total} slides) and {OUT_MD}")
